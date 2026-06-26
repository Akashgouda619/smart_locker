import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def build_presentation():
    # 1. Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Paths to the images
    img_dir = os.path.dirname(os.path.abspath(__file__))
    block_img = os.path.join(img_dir, "block_diagram_1781066393227.png")
    circuit_img = os.path.join(img_dir, "circuit_corrected_1781067987373.png")
    flowchart_img = os.path.join(img_dir, "software_flowchart_1781066433006.png")
    
    # 2. Slide 1: Title Slide (Dark Theme)
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0) # Black
    
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "IoT-Based Smart Rental Locker System"
    p1.font.name = "Arial"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255) # White
    p1.space_after = Pt(12)
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf1.add_paragraph()
    p2.text = "Prototype Design: Integrated Hardware, Web Backend, & UPI Verification Flow"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(163, 163, 163) # Gray
    p2.alignment = PP_ALIGN.LEFT

    # Helper function to add slides with standard layout
    def add_slide_template(title_text):
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(12, 12, 14) # Very dark charcoal #0c0c0e
        
        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        return slide

    # 3. Slide 2: Project Overview & Core Features
    slide2 = add_slide_template("Project Overview & Core Features")
    left_text = (
        "• State-Driven Architecture: Avoids recursive loops and display flickering.\n"
        "• HTTP-Based Polling Loop: The ESP32 polls the server status endpoint every 3 seconds.\n"
        "• Out-of-Band Physical OTP: Dynamic 6-digit retrieval code is shown on the locker screen.\n"
        "• Local SQLite Ledger: SQLite tracks bookings, locker states, and payments dynamically.\n"
        "• verification Page: Web app verifies payment reference number (UTR) to activate storage locks."
    )
    tx_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf2 = tx_box.text_frame
    tf2.word_wrap = True
    for line in left_text.split('\n'):
        p = tf2.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(229, 229, 229)
        p.space_after = Pt(12)

    # 4. Slide 3: Complete System Block Diagram
    slide3 = add_slide_template("1. Complete System Block Diagram")
    if os.path.exists(block_img):
        # Insert image on the left
        slide3.shapes.add_picture(block_img, Inches(0.5), Inches(1.5), width=Inches(6.0), height=Inches(5.0))
    
    # Insert explanation text on the right
    right_box = slide3.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.2))
    tf3 = right_box.text_frame
    tf3.word_wrap = True
    desc_paras = [
        "System Architecture & Interactions:",
        "The prototype coordinates physical hardware locks with web databases and UPI payment confirmation via local network protocols.",
        "1. The ESP32 acts as the physical client, continuously polling the Flask backend via secure Wi-Fi HTTP requests to check current states.",
        "2. The Flask Backend manages relational logic and booking timers, reading and writing to an SQLite transaction database.",
        "3. Users interact via a mobile web portal to book lockers, view instructions, and submit transaction UTR codes to unlock doors.",
        "4. This decoupled system eliminates cellular SMS gateway dependencies by rendering the retrieval OTP directly on the locker's physical TFT screen."
    ]
    for i, para in enumerate(desc_paras):
        p = tf3.add_paragraph() if i > 0 else tf3.paragraphs[0]
        p.text = para
        p.font.name = "Arial"
        p.font.size = Pt(15 if i > 0 else 18)
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(34, 197, 94) # Green title
        else:
            p.font.color.rgb = RGBColor(229, 229, 229)
        p.space_after = Pt(10)

    # 5. Slide 4: Full Circuit Design & Schematic
    slide4 = add_slide_template("2. Full Circuit Design & Connections")
    if os.path.exists(circuit_img):
        slide4.shapes.add_picture(circuit_img, Inches(0.5), Inches(1.5), width=Inches(6.0), height=Inches(5.0))
        
    right_box4 = slide4.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.2))
    tf4 = right_box4.text_frame
    tf4.word_wrap = True
    circuit_text = [
        "Hardware Connections & Interfaces:",
        "The circuit integrates an ESP32 board, a 2.4\" SPI TFT LCD screen, an SG90 analog servo motor, and a DS3231 RTC module.",
        "• ESP32 NodeMCU LX6: Operates at 240 MHz, handles SPI display drawing, I2C timekeeping, and PWM servo control.",
        "• SPI TFT Display: Pin connections CS (GPIO 15), RESET (GPIO 4), D/C (GPIO 2), MOSI (GPIO 23), and SCK (GPIO 18).",
        "• SG90 Servo Control: Powered from Vin (5V) to prevent brown-out resets, driven by a 50Hz PWM signal from GPIO 13.",
        "• DS3231 RTC Module: I2C connections SDA (GPIO 21) and SCL (GPIO 22) for battery-backed timekeeping.",
        "• Common Ground: All components share a common ground (GND) to maintain electrical logic levels."
    ]
    for i, line in enumerate(circuit_text):
        p = tf4.add_paragraph() if i > 0 else tf4.paragraphs[0]
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(15 if i > 0 else 18)
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(34, 197, 94)
        else:
            p.font.color.rgb = RGBColor(229, 229, 229)
        p.space_after = Pt(10)

    # 6. Slide 5: Hardware Datasheet Specs
    slide5 = add_slide_template("Hardware Datasheet Specifications")
    spec_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf5 = spec_box.text_frame
    tf5.word_wrap = True
    specs = [
        "• ESP32 Microcontroller: 3.3V logic level. Drives Servo (GPIO 13), RTC I2C (GPIO 21, 22), and TFT SPI (GPIO 15, 4, 2, 23, 18).",
        "• ILI9341 2.4\" TFT: Screen requires ~3.3V power, 320x240 pixels (RGB), 16-bit color. Draw logic driven by hardware SPI lines.",
        "• SG90 Micro Servo: Stall torque 1.8 kg-cm. Operating voltage 4.8V - 6V. Requires Vin connection since 3.3V rail lacks sufficient surge current. Driven by GPIO 13.",
        "• DS3231 RTC Module: Timekeeping accuracy ±2ppm. Connected via hardware I2C (SDA=GPIO 21, SCL=GPIO 22) at 3.3V.",
        "• Backlight Power: TFT backlighting LED requires 20mA. Driven from 3.3V rail through a 100-ohm current-limiting resistor.",
        "• Common Grounding: Critical for signal integrity between ESP32 (3.3V), RTC (3.3V), TFT (3.3V), and Servo (5V Vin)."
    ]
    for i, line in enumerate(specs):
        p = tf5.add_paragraph() if i > 0 else tf5.paragraphs[0]
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(229, 229, 229)
        p.space_after = Pt(12)

    # 7. Slide 6: Software Flowchart & Methodology
    slide6 = add_slide_template("3. Software State Machine Flowchart")
    if os.path.exists(flowchart_img):
        slide6.shapes.add_picture(flowchart_img, Inches(0.5), Inches(1.5), width=Inches(6.0), height=Inches(5.0))
        
    right_box6 = slide6.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.2))
    tf6 = right_box6.text_frame
    tf6.word_wrap = True
    states_text = [
        "State Transition Matrix:",
        "The firmware functions as a state machine governed by the backend DB booking state:",
        "1. STATE_AVAILABLE: Displays website QR code on screen. Servo is 0° (locked).",
        "2. STATE_PENDING_PAYMENT: Displays payment UPI QR code on screen.",
        "3. STATE_WAITING_FOR_CLOSE: Payment verified, servo unlocks to 90° (door open).",
        "4. STATE_ACTIVE_RENTAL: Locker door locked, TFT screen renders RTC rental timer.",
        "5. STATE_OTP_GENERATED: Renders large centered white card containing 6-digit OTP.",
        "6. STATE_RETRIEVAL_APPROVED: OTP verified, servo unlocks to 90° for item retrieval."
    ]
    for i, line in enumerate(states_text):
        p = tf6.add_paragraph() if i > 0 else tf6.paragraphs[0]
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(15 if i > 0 else 18)
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(34, 197, 94)
        else:
            p.font.color.rgb = RGBColor(229, 229, 229)
        p.space_after = Pt(8)

    # 8. Slide 7: Presentation Q&A and Key Benefits
    slide7 = add_slide_template("Presentation Q&A & Key Benefits")
    qa_box = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf7 = qa_box.text_frame
    tf7.word_wrap = True
    qa_items = [
        "Q: Why use a TFT Screen to display the retrieval OTP?",
        "A: SMS gateways charge recurring fees and have delivery latency. TFT display OTP acts as a zero-cost physical proof of presence factor.",
        "Q: How is the locker servo secured during database timeouts or crashes?",
        "A: Servo requires continuous PWM signals to move. Without power or on crash, it holds position. Status is restored upon ESP32 reboot.",
        "Q: How does automatic payment detection work in the prototype?",
        "A: Users submit their transaction UTR on the web portal. The backend verifies the UTR formatting and sets status to paid."
    ]
    for i, line in enumerate(qa_items):
        p = tf7.add_paragraph() if i > 0 else tf7.paragraphs[0]
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(16)
        if i % 2 == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(34, 197, 94) # Green for questions
        else:
            p.font.color.rgb = RGBColor(229, 229, 229) # White for answers
        p.space_after = Pt(12)

    # Save presentation
    ppt_path = os.path.join(img_dir, "smart_locker_presentation.pptx")
    try:
        prs.save(ppt_path)
        print(f"Presentation saved to: {ppt_path}")
    except PermissionError:
        v2_path = os.path.join(img_dir, "smart_locker_presentation_v2.pptx")
        prs.save(v2_path)
        print(f"Permission denied on original file (likely open in PowerPoint). Saved to: {v2_path}")

if __name__ == "__main__":
    build_presentation()
